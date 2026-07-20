"""Extract plain text from uploaded files, dispatched by MIME type.

Text formats (txt/markdown/csv/json/xml) need no dependencies. Binary office
formats (pdf/docx/pptx/xlsx) use their libraries via lazy imports, so a missing
library degrades that one type to ``unsupported`` rather than breaking the app.

Security: extracted text is treated as *untrusted data* downstream — any
instructions inside a document are never executed, only summarized.
"""

from __future__ import annotations

import csv
import io
import logging
from dataclasses import dataclass

from app.models.enums import SourceStatus

logger = logging.getLogger("app.documents")

_TEXT_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/tab-separated-values",
    "application/json",
    "application/xml",
    "text/xml",
    "text/html",
}


@dataclass(frozen=True)
class ExtractionResult:
    text: str
    status: str  # SourceStatus value
    error: str | None = None


def extract_text(data: bytes, mime_type: str | None, filename: str = "") -> ExtractionResult:
    mime = (mime_type or "").split(";")[0].strip().lower()
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        if mime in _TEXT_TYPES or ext in {"txt", "md", "markdown", "json", "xml", "html", "htm"}:
            return _ok(_decode(data))
        if mime == "text/csv" or ext in {"csv", "tsv"}:
            return _ok(_csv(data))
        if mime == "application/pdf" or ext == "pdf":
            return _ok(_pdf(data))
        if ext == "docx" or "wordprocessingml" in mime:
            return _ok(_docx(data))
        if ext == "pptx" or "presentationml" in mime:
            return _ok(_pptx(data))
        if ext in {"xlsx", "xlsm"} or "spreadsheetml" in mime:
            return _ok(_xlsx(data))
        if mime.startswith("image/"):
            # Image OCR/vision extraction is a later enhancement.
            return ExtractionResult(
                "", SourceStatus.unsupported.value, "image extraction not enabled"
            )
        # Last resort: try to decode as text if it looks textual.
        decoded = _decode(data)
        if decoded.strip():
            return _ok(decoded)
        return ExtractionResult(
            "", SourceStatus.unsupported.value, f"unsupported type: {mime or ext}"
        )
    except _MissingDep as exc:
        return ExtractionResult("", SourceStatus.unsupported.value, str(exc))
    except Exception as exc:  # noqa: BLE001 - one bad file must not break a build
        logger.warning("Extraction failed for %s (%s): %s", filename, mime, exc)
        return ExtractionResult("", SourceStatus.failed.value, str(exc)[:300])


def _ok(text: str) -> ExtractionResult:
    return ExtractionResult(text.strip(), SourceStatus.extracted.value)


class _MissingDep(Exception):
    pass


def _decode(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")


def _csv(data: bytes) -> str:
    rows = csv.reader(io.StringIO(_decode(data)))
    return "\n".join(", ".join(cell for cell in row) for row in rows)


def _pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise _MissingDep("pypdf not installed") from exc
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(data: bytes) -> str:
    try:
        import docx
    except ImportError as exc:
        raise _MissingDep("python-docx not installed") from exc
    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _MissingDep("python-pptx not installed") from exc
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _MissingDep("openpyxl not installed") from exc
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(", ".join(cells))
    return "\n".join(parts)
